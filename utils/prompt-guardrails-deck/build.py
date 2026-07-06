# -*- coding: utf-8 -*-
"""
Build: The Art of Prompt Guardrails - AI Security Maturity (Crawl to Run)
Restyled onto the AccuKnox master template (utils/doc-ppt-template/PPT Template.pptx).
All original slide content is preserved; layout/colors/fonts are rebranded.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = r"D:\Atharva\AccuKnox\HelpDocs\utils\doc-ppt-template\PPT Template.pptx"
OUT = os.path.join(HERE, "Art_of_Prompt_Guardrails_CrawlWalkRun.pptx")
LOGO = os.path.join(HERE, "logo_white.png")
PIPELINE_IMG = os.path.join(HERE, "pipeline_diagram.png")
ICONS = os.path.join(HERE, "icons")

# ---- Brand colors --------------------------------------------------
NAVY      = RGBColor(0x11, 0x20, 0x6D)   # primary brand navy (header bars, badges)
DEEP      = RGBColor(0x00, 0x00, 0xA0)   # deep navy (closing slide bg)
PURPLE    = RGBColor(0x4D, 0x4D, 0xD9)   # accent purple-blue
PURPLE_LT = RGBColor(0x5C, 0x5C, 0xFF)
RED       = RGBColor(0xC8, 0x00, 0x19)   # problem / threat framing
GREEN     = RGBColor(0x16, 0xA5, 0x5C)   # advantage / success
GREEN_DK  = RGBColor(0x0B, 0x7A, 0x42)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x0D, 0x1B, 0x4B)   # body text on white
MID       = RGBColor(0x59, 0x59, 0x59)   # captions / secondary text
FILL_BLUE = RGBColor(0xEE, 0xF3, 0xFF)   # light navy tint card fill
FILL_ALT  = RGBColor(0xF5, 0xF6, 0xFF)   # alternating light tint
FILL_RED  = RGBColor(0xFF, 0xF0, 0xF2)   # light red tint (problem cards)
FILL_GRN  = RGBColor(0xE7, 0xF6, 0xEE)   # light green tint (advantage cards)
BORDER    = RGBColor(0xC4, 0xCC, 0xDE)

FONT = "Montserrat"

# ---- low-level helpers ----------------------------------------------
def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)

def add_slide(prs, layout_name):
    return prs.slides.add_slide(get_layout(prs, layout_name))

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def set_line(shape, color, w=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            sp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        set_fill(sp, fill)
    if line is None:
        no_line(sp)
    else:
        set_line(sp, line, line_w)
    sp.shadow.inherit = shadow
    return sp

def oval(slide, x, y, d, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    set_fill(sp, fill)
    if line is None:
        no_line(sp)
    else:
        set_line(sp, line, 1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def para(tf, text, size=10, bold=False, italic=False, color=DARK, align=PP_ALIGN.LEFT,
         font=FONT, space_before=0, space_after=2, first=False, line_spacing=None, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    prefix = "•  " if bullet else ""
    r = p.add_run()
    r.text = prefix + text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return p

def simple_text(slide, x, y, w, h, text, size=10, bold=False, italic=False, color=DARK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=None, wrap=True):
    tb, tf = textbox(slide, x, y, w, h, anchor=anchor, wrap=wrap)
    para(tf, text, size=size, bold=bold, italic=italic, color=color, align=align, font=font,
         first=True, line_spacing=line_spacing)
    return tb

def add_icon(slide, name, x, y, size=0.32):
    path = os.path.join(ICONS, name + ".png")
    return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(size), height=Inches(size))

def badge(slide, x, y, d, text, fill=NAVY, txtcolor=WHITE, size=13):
    sp = oval(slide, x, y, d, fill)
    tf = sp.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txtcolor; r.font.name = FONT
    return sp

# ---- slide chrome: navy header bar + logo + optional subtitle -------
def add_header(slide, title, subtitle=None, title_size=20, emoji_ok=True):
    bar = rect(slide, -0.05, -0.03, 10.1, 0.80, fill=NAVY)
    logo_w = 1.30
    slide.shapes.add_picture(LOGO, Inches(8.25), Inches(0.20), width=Inches(logo_w))
    title_w = 7.9
    tb, tf = textbox(slide, 0.40, 0.0, title_w, 0.80, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, size=title_size, bold=True, color=WHITE, font=FONT, first=True)
    if subtitle:
        simple_text(slide, 0.40, 0.85, 9.2, 0.32, subtitle, size=11.5, italic=True,
                    color=MID, font=FONT, anchor=MSO_ANCHOR.MIDDLE)
        return 1.24
    return 0.95

def arrow_connector(slide, x, y, w=0.20, h=0.20, color=NAVY):
    sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(sp, color)
    no_line(sp)
    sp.shadow.inherit = False
    return sp

# =====================================================================
prs = Presentation(SRC)
# Drop the 11 example slides that ship with the blank template (and their rels).
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(sldId)

print("Template loaded, layouts cleared. Building content...")

# =====================================================================
# SLIDE 1 -- Title (hero cover: navy bg + product montage + badges all
# inherited from the "AccuKnox Intro Title" layout)
# =====================================================================
s1 = add_slide(prs, "AccuKnox Intro Title")
s1.placeholders[0].text_frame.text = "The Art of Prompt Guardrails"
for p in s1.placeholders[0].text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
s1.placeholders[1].text_frame.text = "AI Security Maturity: From Crawl to Run | AccuKnox & The Broader Ecosystem"
for p in s1.placeholders[1].text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(13); r.font.bold = False; r.font.italic = True
        r.font.color.rgb = RGBColor(0xC9, 0xD3, 0xF5); r.font.name = FONT

# =====================================================================
# SLIDE 2 -- The Evolution of AI
# =====================================================================
s2 = add_slide(prs, "Basic No Content")
add_header(s2, "The Evolution of AI",
           "From attention mechanisms to autonomous, tool-using agents")

evo_cards = [
    ("1", "Transformers", "Attention-based architecture (2017) — the foundation every later layer scales"),
    ("2", "LLMs | RAG | Multimodal", "Language generation + retrieval grounding + vision / audio / text"),
    ("3", "AI Agents | Agentic AI | MCP", "Tool use, autonomy & protocol-based context sharing across systems"),
    ("4", "Agentic AI Architecture", "A new approach: prompts, context, loops and execution harness"),
]
card_w, gap, x0, y0, card_h = 2.05, 0.30, 0.45, 1.35, 1.60
for i, (num, label, desc) in enumerate(evo_cards):
    x = x0 + i * (card_w + gap)
    rect(s2, x, y0, card_w, card_h, fill=FILL_BLUE, line=BORDER, line_w=0.75, rounded=True)
    badge(s2, x + 0.14, y0 + 0.14, 0.36, num, fill=NAVY, size=12)
    simple_text(s2, x + 0.12, y0 + 0.58, card_w - 0.24, 0.42, label, size=11.5, bold=True, color=NAVY)
    simple_text(s2, x + 0.12, y0 + 1.00, card_w - 0.24, card_h - 1.05, desc, size=9, color=DARK, line_spacing=1.05)
    if i < 3:
        arrow_connector(s2, x + card_w + gap * 0.30, y0 + card_h / 2 - 0.09, 0.18, 0.18, color=NAVY)

simple_text(s2, 0.45, 3.08, 9.10, 0.25, "Agentic AI architecture approach", size=11, bold=True, color=NAVY)

pills = ["Prompts", "Context Window", "Loops", "Harness"]
pill_w, pill_gap, py, pill_h = 1.95, 0.42, 3.45, 0.58
for i, label in enumerate(pills):
    x = x0 + i * (pill_w + pill_gap)
    rect(s2, x, py, pill_w, pill_h, fill=WHITE, line=NAVY, line_w=1.25, rounded=True)
    simple_text(s2, x, py, pill_w, pill_h, label, size=11, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        arrow_connector(s2, x + pill_w + pill_gap * 0.30, py + pill_h / 2 - 0.09, 0.18, 0.18, color=PURPLE)

rect(s2, 0.45, 4.55, 9.10, 0.55, fill=FILL_RED, rounded=True)
simple_text(s2, 0.65, 4.55, 8.70, 0.55, "Every new layer of capability adds a new layer of attack surface.",
            size=12.5, bold=True, italic=True, color=RED, anchor=MSO_ANCHOR.MIDDLE)
