from pptx import Presentation

def inch(v):
    return round(v/914400, 2) if v is not None else None

SRC = r"D:\Atharva\AccuKnox\HelpDocs\utils\PPT TEMPLATE - ALWAYS WHEN ASKED TO MAKE PPTS USE THIS.pptx"
p = Presentation(SRC)
slides = list(p.slides)

print("=== SLIDE 1 shapes (all) ===")
for j, sh in enumerate(slides[0].shapes):
    st = str(sh.shape_type).split()[0]
    extra = ""
    if sh.shape_type == 13:
        extra = f" name={sh.name!r}"
    print(f"  [{j}] {st} x={inch(sh.left)} y={inch(sh.top)} w={inch(sh.width)} h={inch(sh.height)}{extra}")

def fill_hex(sh):
    try:
        return str(sh.fill.fore_color.rgb)
    except Exception:
        return None

print("\n=== SLIDE 6 header fills ===")
for idx in [2,3,4,5]:
    print(f"  [{idx}] fill={fill_hex(slides[5].shapes[idx])}")

print("\n=== SLIDE 12 legend + ASPM dots ===")
for idx in [168,170,172,173,174,175,178,200,198,176,179,180,177,181,182,183,184,185,186,187,188,190,191,192,193,194,195]:
    try:
        print(f"  [{idx}] y={inch(slides[11].shapes[idx].top)} fill={fill_hex(slides[11].shapes[idx])}")
    except Exception as e:
        print(f"  [{idx}] err {e}")

print("\n=== SLIDE 1 title placeholder run structure ===")
for idx in [0,1]:
    sh = slides[0].shapes[idx]
    print(f"  [{idx}] paras={len(sh.text_frame.paragraphs)}")
    for pi, para in enumerate(sh.text_frame.paragraphs):
        print(f"     para{pi} runs={len(para.runs)} text={para.text!r}")
        for r in para.runs:
            print(f"        run name={r.font.name} sz={r.font.size.pt if r.font.size else None} bold={r.font.bold}")
