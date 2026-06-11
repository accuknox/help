from pptx import Presentation
from pptx.util import Pt

SRC = r"D:\Atharva\AccuKnox\HelpDocs\utils\PPT TEMPLATE - ALWAYS WHEN ASKED TO MAKE PPTS USE THIS.pptx"
p = Presentation(SRC)
slides = list(p.slides)

def fill_hex(sh):
    try:
        f = sh.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            return str(f.fore_color.rgb)
    except Exception as e:
        return f"err:{e}"
    return None

def first_run(sh):
    try:
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                col = None
                try:
                    if r.font.color and r.font.color.type is not None:
                        col = str(r.font.color.rgb)
                except Exception:
                    col = "theme?"
                sz = r.font.size.pt if r.font.size else None
                return (repr(r.text[:25]), r.font.name, sz, r.font.bold, col)
    except Exception as e:
        return ("err", str(e), None, None, None)
    return None

probe = {
    "S6 before-hdr-red [3]": (6, 3), "S6 after-hdr-navy [5]": (6, 5),
    "S6 after-accentbar [15]": (6, 15), "S6 before-row [6]": (6, 6),
    "S6 after-row [14]": (6, 14), "S6 summary [26]": (6, 26),
    "S7 numbox [6]": (7, 6), "S7 title-chevron [5]": (7, 5), "S7 bullets [7]": (7, 7),
    "S8 statpanel [1]": (8, 1), "S8 statnum [2]": (8, 2), "S8 ucnum-box [14]": (8, 15),
    "S9 card [4]": (9, 4), "S9 accentbar [5]": (9, 5), "S9 cardhdr [6]": (9, 6),
    "S9 statcard [16]": (9, 16), "S9 statnum [17]": (9, 17),
    "S9 impactbar-label [26]": (9, 26), "S9 impactbar-body [25]": (9, 25),
    "S10 statcard [8]": (10, 8), "S10 statnum [9]": (10, 9),
    "S2 numbox [2]": (2, 2), "S2 rowtitle [3]": (2, 3),
    "S1 title [0]": (1, 0),
}
for label, (sl, idx) in probe.items():
    sh = slides[sl-1].shapes[idx]
    print(f"{label}: fill={fill_hex(sh)} run={first_run(sh)}")

# table colors (slide 15)
print("\n--- SLIDE 15 TABLE ---")
tbl = slides[14].shapes[2].table
for ri in range(min(2, len(tbl.rows))):
    for ci in range(min(2, len(tbl.columns))):
        c = tbl.cell(ri, ci)
        fc = None
        try:
            fc = str(c.fill.fore_color.rgb)
        except Exception as e:
            fc = f"err"
        run = None
        for para in c.text_frame.paragraphs:
            for r in para.runs:
                rc = None
                try:
                    rc = str(r.font.color.rgb)
                except Exception:
                    rc = "?"
                run = (repr(r.text[:20]), r.font.name, r.font.size.pt if r.font.size else None, rc)
                break
            if run: break
        print(f"  cell({ri},{ci}) fill={fc} run={run}")
# last AccuKnox col cell with a check
last = len(tbl.columns)-1
c = tbl.cell(1, last)
print(f"  AccuKnox cell(1,{last}) text={c.text_frame.text[:30]!r}")
for para in c.text_frame.paragraphs:
    for r in para.runs:
        try: print("    run", repr(r.text[:20]), "color", str(r.font.color.rgb))
        except Exception: print("    run", repr(r.text[:20]), "color ?")
