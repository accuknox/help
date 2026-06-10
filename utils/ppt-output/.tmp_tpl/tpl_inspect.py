from pptx import Presentation

def inch(v):
    try:
        return round(v / 914400, 2)
    except Exception:
        return None

SRC = r"D:\Atharva\AccuKnox\HelpDocs\utils\PPT TEMPLATE - ALWAYS WHEN ASKED TO MAKE PPTS USE THIS.pptx"
p = Presentation(SRC)
out = []
for i, s in enumerate(p.slides):
    out.append(f"\n========== SLIDE {i+1} (layout={s.slide_layout.name}) ==========")
    for j, sh in enumerate(s.shapes):
        st = str(sh.shape_type).split()[0] if sh.shape_type is not None else "None"
        pos = f"x={inch(sh.left)} y={inch(sh.top)} w={inch(sh.width)} h={inch(sh.height)}"
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text.strip().replace("\n", " / ")
            out.append(f"  [{j}] {st} {pos} TXT='{txt[:140]}'")
        elif sh.shape_type == 13:
            continue
        elif sh.has_table:
            out.append(f"  [{j}] TABLE {pos} rows={len(sh.table.rows)} cols={len(sh.table.columns)}")
        else:
            out.append(f"  [{j}] {st} {pos}")
    npic = sum(1 for sh in s.shapes if sh.shape_type == 13)
    out.append(f"  (+{npic} pictures)")
open(r"D:\Atharva\AccuKnox\HelpDocs\.tmp_tpl\shapes.txt", "w", encoding="utf-8").write("\n".join(out))
print("WROTE shapes.txt, lines:", len(out))
