# -*- coding: utf-8 -*-
from pptx import Presentation
NL = chr(10)
OUTPUT_PPTX = r"D:\Atharva\AccuKnox\HelpDocs\utils\ppt-output\<replace-with-output-filename>.pptx"
p = Presentation(OUTPUT_PPTX)
slides = list(p.slides)
print("SLIDES:", len(slides))
bad = 0
for i, s in enumerate(slides):
    title = ""
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            title = sh.text_frame.text.strip().split(NL)[0][:55]
            break
    npic = sum(1 for sh in s.shapes if sh.shape_type == 13)
    pics = f"  [{npic} img]" if npic else ""
    print(f"  {i+1:2d}. {title}{pics}")
    for sh in s.shapes:
        if sh.has_text_frame and ("—" in sh.text_frame.text or "–" in sh.text_frame.text):
            bad += 1
print("EM/EN DASH HITS:", bad)
