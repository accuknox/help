# -*- coding: utf-8 -*-
from pptx import Presentation
OUT = r"D:\Atharva\AccuKnox\HelpDocs\utils\ppt-output\<replace-with-output-filename>.pptx"
prs = Presentation(OUT)
bad = []
def scan_tf(tf, where):
    for para in tf.paragraphs:
        t = "".join(r.text for r in para.runs)
        if "—" in t or "–" in t:
            bad.append((where, t))
for i, s in enumerate(prs.slides):
    for sh in s.shapes:
        if sh.has_text_frame:
            scan_tf(sh.text_frame, f"slide{i+1}")
        if sh.has_table:
            for ri, row in enumerate(sh.table.rows):
                for ci, cell in enumerate(row.cells):
                    scan_tf(cell.text_frame, f"slide{i+1} tbl[{ri},{ci}]")
print("EM/EN DASH HITS:", len(bad))
for w, t in bad:
    print("  ", w, "::", t[:90])
print("TOTAL SLIDES:", len(list(prs.slides)))
