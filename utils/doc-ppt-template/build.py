# -*- coding: utf-8 -*-
"""
Build: Bandhan AMC — Application Security Platform Proposal
Derived from the AccuKnox master template (utils/PPT TEMPLATE...).
Edits the real template in place to preserve all brand assets, then
restructures into a phased SAST proposal deck.
"""
import copy, shutil, os, glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image as PILImage

SRC   = r"D:\Atharva\AccuKnox\HelpDocs\utils\PPT TEMPLATE - ALWAYS WHEN ASKED TO MAKE PPTS USE THIS.pptx"
OUT   = r"D:\Atharva\AccuKnox\HelpDocs\utils\ppt-output\Bandhan_AMC_Application_Security_Proposal.pptx"
BLANK = r"D:\Atharva\AccuKnox\HelpDocs\utils\ppt-output\AccuKnox_Proposal_Template_BLANK.pptx"

# ---- Brand colors -------------------------------------------------
NAVY      = RGBColor(0x11,0x20,0x6D)
NAVY_TTL  = RGBColor(0x0E,0x35,0x94)
PURPLE    = RGBColor(0x4D,0x4D,0xD9)
RED       = RGBColor(0xC8,0x00,0x19)
WHITE     = RGBColor(0xFF,0xFF,0xFF)
GREEN     = RGBColor(0x16,0xA5,0x5C)   # AK accent green (live / advantage)
GREEN_DK  = RGBColor(0x0B,0x7A,0x42)
GREEN_LT  = RGBColor(0xE7,0xF6,0xEE)
GREY_PH   = RGBColor(0xEE,0xF0,0xF6)   # screenshot placeholder fill
GREY_BD   = RGBColor(0xC4,0xCC,0xDE)   # placeholder border
LAV       = RGBColor(0xE6,0xE6,0xFA)
NOT_ON    = RGBColor(0x99,0x00,0xFF)
FONT      = "Space Grotesk"

# ---- helpers ------------------------------------------------------
def settext(shape, *lines):
    """Replace text of non-empty paragraphs in order, preserving formatting."""
    tf = shape.text_frame
    idx = 0
    paras = list(tf.paragraphs)
    for para in paras:
        if para.runs and para.text.strip():
            if idx < len(lines):
                para.runs[0].text = lines[idx]
                for r in para.runs[1:]:
                    r._r.getparent().remove(r._r)
                idx += 1
            else:
                para.runs[0].text = ""
                for r in para.runs[1:]:
                    r._r.getparent().remove(r._r)
    # if shape had a single empty paragraph, force-set first line
    if idx == 0 and lines:
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = lines[0]
        else:
            run = p.add_run(); run.text = lines[0]; run.font.name = FONT

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def run_color(shape, color):
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            r.font.color.rgb = color

def add_box(slide, x, y, w, h, text, size, color=NAVY, bold=True,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT,
            fill=None, line=None, line_w=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    if text:
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return sp

def rounded_badge(slide, x, y, w, h, text, fill, txtcolor=WHITE, size=10):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txtcolor
    return sp

# ===================================================================
shutil.copy(SRC, OUT)
prs = Presentation(OUT)
S = list(prs.slides)
def sh(slide_i, shape_i):   # 1-based slide, 0-based shape
    return S[slide_i-1].shapes[shape_i]

# ---- Phase 3 screenshots ------------------------------------------
ASSETS = r"D:\Atharva\AccuKnox\HelpDocs\proposal-output\assets"
os.makedirs(ASSETS, exist_ok=True)
_shots = sorted(glob.glob(r"D:\Atharva\AccuKnox\HelpDocs\utils\sast-images-phase3-cd*"))
# crop the product-name breadcrumb strip off shot 0
IMG = {i: _shots[i] for i in range(len(_shots))}
if _shots:
    # crop the whole top header band off shot 0 (removes product-name breadcrumb + login email)
    _im0 = PILImage.open(_shots[0]).convert("RGB"); _w0, _h0 = _im0.size
    _crop0 = os.path.join(ASSETS, "shot0_crop.png")
    _im0.crop((0, int(_h0 * 0.225), _w0, _h0)).save(_crop0)
    IMG[0] = _crop0

def img_fit(slide, path, x, y, w, h, caption=None):
    """Place an image scaled to fit (w,h) preserving aspect, centered, framed; optional caption below."""
    im = PILImage.open(path); ar = im.size[0] / im.size[1]
    if w / h > ar:
        dh = h; dw = h * ar
    else:
        dw = w; dh = w / ar
    dx = x + (w - dw) / 2; dy = y + (h - dh) / 2
    pic = slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    pic.line.color.rgb = GREY_BD; pic.line.width = Pt(1.25)
    pic.shadow.inherit = False
    if caption:
        add_box(slide, x, y + h + 0.03, w, 0.3, caption, 9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    return pic

# ---------- SLIDE 1 — Title ----------
t0 = sh(1,0)
settext(t0, "Application Security Platform Proposal")
t0.left = Inches(0.76); t0.top = Inches(2.30); t0.width = Inches(5.45); t0.height = Inches(0.9)
for _p in t0.text_frame.paragraphs:
    for _r in _p.runs:
        _r.font.size = Pt(22)
settext(sh(1,1), "")   # blank the old date placeholder
# remove Danube logo picture [3], add client name + subtitle
danube = sh(1,3)
danube._element.getparent().remove(danube._element)
add_box(S[0], 0.76, 1.12, 5.2, 0.56, "BANDHAN AMC", 29, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_box(S[0], 0.78, 1.73, 5.2, 0.30, "Asset Management Company", 12.5, color=RGBColor(0xB8,0xC4,0xE8), bold=False, align=PP_ALIGN.LEFT)
add_box(S[0], 0.78, 3.26, 5.6, 0.30, "SAST  ·  Secrets Scanning  ·  SBOM", 12, color=RGBColor(0xCE,0xD8,0xF2), bold=True, align=PP_ALIGN.LEFT)
add_box(S[0], 0.78, 3.60, 5.6, 0.30, "Air-Gapped GitLab, On-Prem  ·  June 2026", 11.5, color=RGBColor(0xB8,0xC4,0xE8), bold=False, align=PP_ALIGN.LEFT)

# ---------- SLIDE 2 — Outline (6 rows) ----------
s2 = S[1]
settext(sh(2,0), "Outline")
# existing number boxes [2,4,6,8] and titles [3,5,7,9]
num_boxes = [sh(2,2), sh(2,4), sh(2,6), sh(2,8)]
titles    = [sh(2,3), sh(2,5), sh(2,7), sh(2,9)]
agenda = [
    "AccuKnox Platform Overview",
    "PoC Recap: SAST, Secrets, SBOM",
    "The Phased SAST Roadmap",
    "Phases 1 to 3: On-Prem, AI, Platform",
    "Why AccuKnox vs. Standalone AppSec",
    "Recommended Next Steps",
]
ROW_Y0, ROW_H, ROW_GAP = 1.02, 0.62, 0.10
def place_row(nb, tt, i):
    y = ROW_Y0 + i*(ROW_H+ROW_GAP)
    nb.top = Inches(y); nb.height = Inches(ROW_H)
    tt.top = Inches(y); tt.height = Inches(ROW_H)
# clone two more pairs from box[3]/title[3]  (append order: nb,tt,nb,tt)
for k in range(2):
    nb_el = copy.deepcopy(num_boxes[3]._element)
    tt_el = copy.deepcopy(titles[3]._element)
    s2.shapes._spTree.append(nb_el)
    s2.shapes._spTree.append(tt_el)
new_shapes = list(s2.shapes)[-4:]   # [nb0, tt0, nb1, tt1]
num_boxes.append(new_shapes[0]); titles.append(new_shapes[1])
num_boxes.append(new_shapes[2]); titles.append(new_shapes[3])
for i in range(6):
    place_row(num_boxes[i], titles[i], i)
    settext(num_boxes[i], str(i+1))
    settext(titles[i], agenda[i])

# ---------- SLIDE 3 — About (unchanged) ----------

# ---------- SLIDE 5 — Platform (keep) ----------

# ---------- SLIDE 7 — Current State (PoC recap) ----------
settext(sh(7,0), "Current State: What the PoC Validated")
# Row1 (01) title[17] bullets[19]  (short, single-line bullets)
settext(sh(7,17), "SAST Scanning")
settext(sh(7,19), "Static analysis on GitLab source",
                  "Findings sent to AccuKnox Console", "Status:  Live")
settext(sh(7,11), "Secrets Scanning")
settext(sh(7,13), "Hardcoded credentials & tokens",
                  "Metadata-only upload, no egress", "Status:  Live")
settext(sh(7,5),  "SBOM Generation & Comparison")
settext(sh(7,7),  "CycloneDX / SPDX, generated locally",
                  "Cross-version drift detection", "Status:  Live")
settext(sh(7,23), "Air-Gapped Deployment")
settext(sh(7,25), "Single-node K3s, fully offline",
                  "No public-internet dependency", "Zero source-code egress")
settext(sh(7,26), "PoC completed May 2026   ·   3 / 3 use cases validated")
sh(7,26).left = Inches(4.55); sh(7,26).width = Inches(5.0)
for _r in sh(7,26).text_frame.paragraphs[0].runs:
    _r.font.size = Pt(12)

# ---------- SLIDE 6 — Before / After (SAST) ----------
settext(sh(6,0), "SAST Today  vs.  SAST with AccuKnox")
settext(sh(6,7),  "✗   High false-positive volume drowns real issues")
settext(sh(6,9),  "✗   Findings arrive without remediation guidance")
settext(sh(6,11), "✗   Siloed scanners, no unified triage")
settext(sh(6,13), "✗   Manual severity assessment slows the team")
settext(sh(6,16), "✓   AI triage cuts false positives, sharpens signal")
settext(sh(6,19), "✓   “Ask AI” remediation guidance on any finding")
settext(sh(6,22), "✓   Unified console: SAST, Secrets, SBOM, SCA")
settext(sh(6,25), "✓   EPSS / CISA KEV risk-based prioritization")
settext(sh(6,26), "AccuKnox layers AI-driven quality onto your GitLab SAST, without code leaving your network.")
# green accents on AFTER column (accent bars [15,18,21,24]) + green check color
for bar in [15,18,21,24]:
    set_fill(sh(6,bar), GREEN)
for txt in [16,19,22,25]:
    run_color(sh(6,txt), NAVY)  # keep body navy; check stays inline

# ---------- SLIDE 8 — Phased SAST Roadmap ----------
settext(sh(8,27), "The Phased SAST Roadmap")
settext(sh(8,12), "A Phased Path to Advanced SAST")
settext(sh(8,13), "From validated on-prem SAST to AI-powered, platform-grade application security")
settext(sh(8,17), "Phase 1:  On-Prem SAST   ·   Live Today")
settext(sh(8,21), "Phase 2:  AI-Powered SAST   ·   Near-Term")
settext(sh(8,25), "Phase 3:  Advanced SAST Platform   ·   2-3 Months")
# phase number squares [14(01),18(02),22(03)] -> green / purple / navy progression
set_fill(sh(8,14), GREEN)
set_fill(sh(8,18), PURPLE)
set_fill(sh(8,22), NAVY)
# stat panel
settext(sh(8,2), "3 / 3"); settext(sh(8,3), "PoC use cases live")
settext(sh(8,5), "0");     settext(sh(8,6), "Source-code egress")
settext(sh(8,8), "85%");   settext(sh(8,9), "Target noise reduction (AI)")
settext(sh(8,11), "2-3");  settext(sh(8,26), "Months to platform GA")

# ---------- SLIDE 9 — Phase 1: On-Prem SAST (Live) ----------
settext(sh(9,0), "Phase 1:  On-Prem SAST   ·   Implemented")
settext(sh(9,3), "Validated in the Bandhan AMC PoC: static analysis runs entirely inside your air-gapped GitLab environment, with findings centralized on the on-prem AccuKnox Console.")
settext(sh(9,6),  "Air-Gapped Execution")
settext(sh(9,7),  "The ASPM scanner runs on a Bandhan AMC internal VM using the host Docker daemon, with no privileged Docker-in-Docker and no scanner images leaving the network.")
settext(sh(9,10), "GitLab-Native Onboarding")
settext(sh(9,11), "Source is pulled from your internal GitLab. A .env supplies scan parameters; findings are pushed automatically to the AccuKnox Console.")
settext(sh(9,14), "Centralized, Correlated Findings")
settext(sh(9,15), "SAST results are normalized and unified with Secrets and SBOM findings for one-place triage, trend analysis, and reporting.")
settext(sh(9,17), "100%"); settext(sh(9,18), "On-prem, runs inside your network")
settext(sh(9,20), "0");    settext(sh(9,21), "Source code or artifacts egressing")
settext(sh(9,23), "3 / 3"); settext(sh(9,24), "PoC use cases onboarded & validated")
settext(sh(9,27), "Live Today")
settext(sh(9,28), "SAST, Secrets, and SBOM operate locally, satisfying Bandhan AMC's data-residency and confidentiality requirements.")
# green accents: left card bars [5,9,13], impact label [26] green, stat numbers green
for bar in [5,9,13]:
    set_fill(sh(9,bar), GREEN)
set_fill(sh(9,26), GREEN)
for n in [17,20,23]:
    run_color(sh(9,n), GREEN_DK)

# ---------- SLIDE 10 — Phase 2: AI-Powered SAST ----------
settext(sh(10,0), "Phase 2:  AI-Powered SAST")
settext(sh(10,3), "The near-term enhancement: AI applied to SAST output for measurable quality gains, with no cloud migration required.")
settext(sh(10,9),  "AI Triage"); settext(sh(10,10), "Validates findings & suppresses noise automatically")
settext(sh(10,12), "Ask AI");    settext(sh(10,13), "Guided remediation on any finding, individual or batch")
settext(sh(10,15), "Enrich");    settext(sh(10,16), "Severity, CWE class & actionable summaries auto-added")
settext(sh(10,19), "AI False-Positive Reduction")
settext(sh(10,20), "Cross-scan correlation across SAST, Secrets, SBOM and SCA, with EPSS scoring, CISA KEV, CWE classification and business-impact weighting. AI identifies false positives so teams focus on genuinely exploitable issues.")
settext(sh(10,23), "AI-Assisted Remediation & Enrichment")
settext(sh(10,24), "“Ask AI” (AskADA copilot) gives step-by-step fixes on any finding type, supports batch remediation, and auto-generates severity assessments and concise, actionable summaries.")
settext(sh(10,6), "Why It Matters")
settext(sh(10,7), "AI sharpens signal-to-noise and accelerates fixes, directly addressing the false-positive pain in current SAST output.")
# make the 3 stat numbers smaller-feel by recoloring to purple accent
for n in [9,12,15]:
    run_color(sh(10,n), PURPLE)
set_fill(sh(10,6), PURPLE)

# ---------- SLIDE 11 — Phase 3: Advanced SAST Platform ----------
settext(sh(11,0), "Phase 3:  Advanced SAST Platform")
settext(sh(11,3), "On the roadmap (2-3 months): a platform-grade SAST experience with advanced detection, deeper analytics, and richer dashboards.")
# left 3 cards: number -> 01/02/03, title, body
settext(sh(11,10), "01"); settext(sh(11,11), "Advanced Detection")
settext(sh(11,12), "Deeper rule coverage and contextual analysis across more languages and frameworks.")
settext(sh(11,15), "02"); settext(sh(11,16), "Analytics & Trends")
settext(sh(11,17), "Historical trends, severity movement, and top-vulnerable repos and files at a glance.")
settext(sh(11,20), "03"); settext(sh(11,21), "Richer Dashboards")
settext(sh(11,22), "Executive and engineer views with export-ready ASPM reporting for audit and stakeholders.")
for n in [10,15,20]:
    run_color(sh(11,n), PURPLE)
# right column -> clear SCA cards, drop in a hero product screenshot
_els = [sh(11,_i)._element for _i in [23,24,25,26,27,28,29,30,31,32,33,34,35]]
for _e in _els:
    _e.getparent().remove(_e)
add_box(S[10], 5.25, 1.18, 4.3, 0.3, "Product Preview", 13, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
img_fit(S[10], IMG[3], 5.2, 1.55, 4.4, 2.45, "Code Analysis dashboard:  noise reduced 47%")
settext(sh(11,6), "Timeline")
settext(sh(11,7), "Targeted for general availability within 2-3 months; feature previews available on request.")
set_fill(sh(11,6), PURPLE)

# ---------- SLIDE 12 — Platform breadth ----------
settext(sh(12,0), "AccuKnox Platform:  Unified ASPM & Beyond")
settext(sh(12,171), "Onboarded in Bandhan AMC PoC")
settext(sh(12,169), "Available on platform")
# recolor ASPM dots: only SAST[172], Secret[174], SBOM[178] green; rest not-onboarded
aspm_green = [172,174,178]
aspm_all   = [172,173,174,175,178,200,198]
for idx in aspm_all:
    set_fill(sh(12,idx), GREEN if idx in aspm_green else NOT_ON)
# all other columns -> not onboarded
for idx in [176,179,180,177, 181,182,183,184, 185,186,187,188, 190,191,192,193,194,195]:
    set_fill(sh(12,idx), NOT_ON)
# legend marker green -> AK green
set_fill(sh(12,170), GREEN)

# ---------- SLIDE 13 — Case study (reframe) ----------
settext(sh(13,0), "Customer Proof:  Regulated, On-Prem")
# keep headline/sub/outcomes (real proof point); clean em dashes in retained text
for _para in sh(13,7).text_frame.paragraphs:
    for _r in _para.runs:
        if "—" in _r.text:
            _r.text = _r.text.replace(" — ", ", ").replace("—", ", ")

# ---------- SLIDE 15 — Why AccuKnox (rebuild comparison table) ----------
s15 = S[14]
settext(sh(15,0), "Why AccuKnox:  Advanced AppSec & SAST")
old_tbl = sh(15,2)
old_tbl._element.getparent().remove(old_tbl._element)
rows_data = [
    ("Capability", "Invicti", "AccuKnox"),
    ("Native AI SAST engine", "⚠️  No native SAST (3rd-party Mend)", "✅  Native AI-accelerated SAST + enrichment"),
    ("AI false-positive reduction", "✅  Proof-based scanning (DAST)", "✅  EPSS · CISA KEV · AI FP-ID · cross-scan"),
    ("AI-assisted remediation", "⚠️  DAST findings only", "✅  Ask AI on any finding + batch (AskADA)"),
    ("Secrets · SBOM · SCA with SAST", "⚠️  SBOM only; 3rd-party SCA", "✅  Native Secrets, SBOM/CBOM/AIBOM, SCA"),
    ("Cloud · K8s · Runtime security", "❌  Not available", "✅  CSPM · CWPP · KSPM · CDR"),
    ("Air-gapped / on-prem parity", "⚠️  On-prem (limited)", "✅  Full parity · air-gapped · OEM/MSSP"),
    ("Compliance frameworks", "⚠️  8 frameworks", "✅  33+ frameworks"),
]
nrows, ncols = len(rows_data), 3
tx, ty, tw, th = Inches(0.45), Inches(1.05), Inches(9.1), Inches(3.95)
gtbl = s15.shapes.add_table(nrows, ncols, tx, ty, tw, th).table
gtbl.first_row = False; gtbl.horz_banding = False
gtbl.columns[0].width = Inches(3.0)
gtbl.columns[1].width = Inches(3.05)
gtbl.columns[2].width = Inches(3.05)
for ri, row in enumerate(rows_data):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE if ci < 2 else GREEN
        elif ci == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = LAV
        elif ci == 2:
            cell.fill.solid(); cell.fill.fore_color.rgb = GREEN_LT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        tfc = cell.text_frame; tfc.word_wrap = True
        p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.name = FONT
        r.font.size = Pt(11 if ri == 0 else 9.0)
        r.font.bold = (ri == 0 or ci == 0)
        if ri == 0:
            r.font.color.rgb = WHITE
        elif ci == 0:
            r.font.color.rgb = NAVY
        else:
            r.font.color.rgb = NAVY
# footer note
add_box(s15, 0.45, 5.07, 9.1, 0.33,
        "AccuKnox integrates natively with GitLab CI/CD, adding the AI SAST layer on top of your existing pipeline.",
        9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, fill=GREEN_LT)

# ---------- SLIDE 14 — Recommended Next Steps (add content to section break) ----------
s14 = S[13]
settext(sh(14,0), "Recommended Next Steps")
# move the title up
sh(14,0).top = Inches(0.7)
steps = [
    ("01", "Operationalize the scanners",
     "Integrate SAST, Secrets & SBOM into GitLab CI/CD for automated, shift-left coverage on every push and merge request."),
    ("02", "Enable Phase 2:  AI SAST",
     "Turn on AI triage, false-positive reduction and “Ask AI” remediation to lift finding quality immediately."),
    ("03", "Plan Phase 3:  Platform Rollout",
     "Schedule the advanced SAST platform (2-3 months) and review feature previews with the AccuKnox team."),
]
sy = 1.95
for num, title, body in steps:
    add_box(s14, 0.9, sy, 0.7, 0.7, num, 22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, fill=GREEN)
    add_box(s14, 1.8, sy-0.02, 7.3, 0.36, title, 15, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_box(s14, 1.8, sy+0.32, 7.3, 0.42, body, 11, color=RGBColor(0xCE,0xD8,0xF2), bold=False, align=PP_ALIGN.LEFT)
    sy += 1.0

# ---------- SLIDE 18 — Closing (keep) ----------

# ---------- SLIDES 16 & 17 — Phase 3 Product Preview galleries ----------
def build_gallery(slide, title, intro, big, smalls):
    settext(slide.shapes[0], title)               # title placeholder in header
    for shp in list(slide.shapes):                # drop the old comparison table
        if shp.has_table:
            shp._element.getparent().remove(shp._element)
    add_box(slide, 0.45, 0.72, 9.1, 0.32, intro, 12, color=NAVY, bold=False, align=PP_ALIGN.LEFT)
    bidx, bcap = big                              # one large featured shot, left
    img_fit(slide, IMG[bidx], 0.45, 1.32, 5.4, 2.95, bcap)
    for (idx, cap), y in zip(smalls, [1.32, 3.07]):   # two supporting shots, stacked right
        img_fit(slide, IMG[idx], 6.05, y, 3.5, 1.4, cap)

build_gallery(S[15], "Phase 3:  Product Preview",
              "From scan results to vulnerability detail and AI-assisted remediation.",
              (0, "Repository findings overview"),
              [(1, "Vulnerability detail with CWE mapping"),
               (2, "AI-assisted remediation")])
build_gallery(S[16], "Phase 3:  Product Preview  (2 / 2)",
              "Deep context: data-flow, code-flow and risk-based prioritization.",
              (4, "Data-flow graph (source to sink)"),
              [(5, "Code-flow analysis"),
               (6, "Risk analysis & business impact")])

# ===================================================================
# Delete slide 4 (redundant platform matrix); reorder, weaving galleries after Phase 3.
desired = [1,2,3,5,7,6,8,9,10,11, 16,17, 12,15,13,14,18]
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
keep_els = [ids[i-1] for i in desired]
delete_els = [el for el in ids if el not in keep_els]
# drop relationships for deleted slides
for el in delete_els:
    rId = el.get(qn('r:id'))
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
for el in ids:
    sldIdLst.remove(el)
for el in keep_els:
    sldIdLst.append(el)

prs.save(OUT)
print("Saved proposal:", OUT, "slides:", len(list(prs.slides)))

# ===================================================================
# Also emit a neutral BLANK template (light Danube de-branding).
prs2 = Presentation(SRC)
T = list(prs2.slides)
def t_sh(si, xi): return T[si-1].shapes[xi]
settext(t_sh(1,0), "[ Presentation Title ]")
settext(t_sh(1,1), "[ Date ]")
repl = {
    6:  {0:"[ Client ]:  Before vs After"},
    8:  {27:"[ Client ] - Business Use Cases", 12:"[ Headline ]"},
    9:  {0:"Business Case 01  ·  [ Title ]"},
    10: {0:"Business Case 02  ·  [ Title ]"},
    11: {0:"Business Case 03  ·  [ Title ]"},
    13: {0:"[ Customer ] - Case Study"},
}
for si, m in repl.items():
    for xi, val in m.items():
        try: settext(t_sh(si,xi), val)
        except Exception: pass
prs2.save(BLANK)
print("Saved blank template:", BLANK)
